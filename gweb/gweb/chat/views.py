import json
import logging
import base64
import hashlib
from uuid import UUID, uuid4

from django.conf import settings
from django.contrib.auth.models import User
from django.shortcuts import get_object_or_404
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from django.utils import timezone
from django.db.models import F
from django.db.models.functions import Greatest

from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status, permissions
from rest_framework.permissions import IsAuthenticated
from rest_framework.generics import RetrieveAPIView
from django.http import FileResponse
from rest_framework.parsers import MultiPartParser, FormParser, JSONParser

from openai import OpenAI
from django.http import StreamingHttpResponse

from .models import Chat, ChatHistory, ChatCollaborator
from .serializers import (
    ChatSerializer,
    ChatHistorySerializer,
    ChatListSerializer,
    ChatMessageSerializer,
    ChatCollaboratorSerializer,
)
from django.core.files.base import ContentFile
from django.core.cache import cache
from.models import TextToSpeech
from .utils import get_client_ip, generate_chat_title_from_openai, get_openai_client

from g_auth.models import GuestChatTracker, GuestIPTracker
from payments.models import UserProfile as PaymentsUserProfile
from files.models import File as StoredFile
import re

logger = logging.getLogger(__name__)
FREE_AUDIO_MINUTES_LIMIT = 10.0
TTS_PERSIST_AUDIO_FILES = bool(getattr(settings, "TTS_PERSIST_AUDIO_FILES", False))
TTS_PERSIST_CHAT_HISTORY = bool(getattr(settings, "TTS_PERSIST_CHAT_HISTORY", False))
TTS_RETURN_INLINE_AUDIO = bool(getattr(settings, "TTS_RETURN_INLINE_AUDIO", True))

STUDY_METHOD_PROMPTS = {
    "feynman": "Explain ideas in simple language first, then expose misunderstandings or gaps and rebuild the concept clearly.",
    "active_recall": "Turn parts of the lesson into short memory-check questions before revealing the answer.",
    "spaced_repetition": "Suggest what to review later and highlight the facts worth revisiting over time.",
    "socratic": "Use guided questions that help the learner reason toward the answer instead of only stating conclusions.",
    "interleaving": "Contrast related ideas, compare cases, and mix connected concepts when it improves understanding.",
    "exam_drill": "Frame the explanation in an exam-focused way with likely question angles, marking cues, or quick practice checks.",
}

def _sanitize_response_text(text: str) -> str:
    if not isinstance(text, str):
        return text
    cleaned = text

    # Preserve readable structure: keep line breaks, normalize spacing per line.
    lines = []
    for line in cleaned.splitlines():
        compact = re.sub(r"[ \t]+", " ", line).strip()
        lines.append(compact)
    cleaned = "\n".join(lines)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    return cleaned


def _sse_payload(payload):
    return f"data: {json.dumps(payload)}\n\n"


def _parse_study_methods(value):
    if value is None or value == "":
        return []

    raw_items = value
    if isinstance(value, str):
        parsed = None
        try:
            parsed = json.loads(value)
        except (TypeError, ValueError, json.JSONDecodeError):
            parsed = [item.strip() for item in value.split(",")]
        raw_items = parsed

    if not isinstance(raw_items, list):
        raw_items = [raw_items]

    cleaned = []
    seen = set()
    for item in raw_items:
        normalized = str(item or "").strip().lower()
        if not normalized or normalized not in STUDY_METHOD_PROMPTS or normalized in seen:
            continue
        seen.add(normalized)
        cleaned.append(normalized)
    return cleaned


def _normalize_study_custom_prompt(value):
    text = str(value or "").strip()
    if not text:
        return ""
    text = re.sub(r"\s+", " ", text)
    return text[:240]


def _build_chat_system_prompt(study_methods=None, study_custom_prompt=""):
    study_methods = study_methods or []
    study_guidance_lines = []
    for method in study_methods:
        guidance = STUDY_METHOD_PROMPTS.get(method)
        if guidance:
            study_guidance_lines.append(f"- {guidance}")

    dynamic_section = ""
    if study_guidance_lines or study_custom_prompt:
        parts = ["\n## Active Study Mode"]
        parts.append("- Apply the following response preferences unless the user asks for a different format.")
        parts.extend(study_guidance_lines)
        if study_custom_prompt:
            parts.append(f"- Additional user instruction: {study_custom_prompt}")
        dynamic_section = "\n".join(parts) + "\n"

    return rf"""You are QwikStudi, an AI study buddy created by Glinax Tech Innovations.

## Your Core Identity
- You are a helpful, patient, and encouraging study companion.
- You help users learn and understand concepts deeply, not just memorize answers.
- When asked for your creator, always answer "Glinax Tech Innovations".
- When asked for your name, always answer "QwikStudi".
- Never say you have no name. Never claim to be created by OpenAI or any other organization.
- Stay in character as a friendly tutor at all times.

## Teaching Approach
- Prioritize understanding: explain the "why" behind each step.
- Use active learning: guide users to discover solutions, not only receive results.
- Use concrete examples before abstractions when possible.
- Offer an alternate explanation when the first one may be too advanced.

## Computational Reasoning Standards (Math, Engineering, CS, Data)
- Start by identifying: what is given, what is unknown, and what method will be used.
- If data is missing or ambiguous, ask one concise clarifying question; if proceeding, state assumptions clearly.
- Show formulas or governing principles first, then substitute values, then compute.
- Keep units attached through calculations and convert units explicitly when needed.
- Prefer exact values when useful, then provide practical rounded results with sensible precision.
- Perform a quick validation check (units, sign, scale, bounds, or reasonableness).
- For code/algorithm tasks, include: approach, complexity (if relevant), and a correct minimal example.
- For multi-step problems, label steps in a clear sequence and end with a direct final answer.

## Presentation Format For Quantitative Problems
- Problem setup
- Given / Find
- Assumptions
- Step-by-step solution
- Final answer (with units)
- Quick check

## Response Style
- Be clear, concise, and student-friendly.
- Break complex work into digestible steps.
- Keep responses scannable with short paragraphs or bullets.
- Match detail level to request complexity and selected response style.
- Use clean Markdown structure (headings, bullets, numbered steps, code fences, and tables when helpful).

## Context Awareness
- Use the current chat context and relevant prior messages.
- Connect new concepts to what the learner already discussed.
{dynamic_section}
## Response Quality Rules
- If uncertain, explicitly state uncertainty instead of guessing.
- Do not invent constants, formulas, citations, experimental values, or standards.
- Prefer readable plain notation for simple expressions (x^2, x_i, x^(n+1), F = m * a).
- For advanced math, you may use Markdown math delimiters: inline $...$ and block $$...$$.
- Do not output raw LaTeX fragments outside math delimiters.
- If you use LaTeX syntax (e.g., \frac, \Phi, \theta, \sum, \int), always place it inside $...$ or $$...$$.
- When solving computational problems, format calculations so each step is visually separable and easy to verify.

Remember: Your goal is to help users become better learners and better problem-solvers."""


def _generate_chat_completion_text(client, model: str, messages: list[dict], temperature: float, max_tokens: int) -> str:
    assembled_parts: list[str] = []
    working_messages = list(messages)

    # One continuation pass if the model hits token limit.
    for _ in range(2):
        response = client.chat.completions.create(
            model=model,
            messages=working_messages,
            temperature=temperature,
            max_tokens=max_tokens,
        )
        choice = response.choices[0]
        content = (choice.message.content or "").strip()
        if content:
            assembled_parts.append(content)

        if choice.finish_reason != "length":
            break

        working_messages.extend(
            [
                {"role": "assistant", "content": content},
                {
                    "role": "user",
                    "content": "Continue exactly where you stopped. Do not repeat prior text. Finish the answer.",
                },
            ]
        )

    return "\n".join(part for part in assembled_parts if part).strip()


def _get_chat_for_write(user, chat_id):
    if not chat_id:
        return None
    try:
        chat = Chat.objects.get(id=chat_id)
    except Chat.DoesNotExist:
        return None
    if chat.user == user:
        return chat
    collaborator = ChatCollaborator.objects.filter(
        chat=chat, collaborator=user, is_approved=True, access_level="edit"
    ).first()
    return chat if collaborator else None

def _get_payments_profile(user):
    profile, _ = PaymentsUserProfile.objects.get_or_create(user=user)
    today = timezone.now().date()
    if profile.is_premium and profile.premium_expiry and profile.premium_expiry < today:
        profile.is_premium = False
        profile.save(update_fields=["is_premium"])
    return profile

def _is_premium(profile):
    today = timezone.now().date()
    return bool(profile.is_premium and profile.premium_expiry and profile.premium_expiry >= today)

def _estimate_tts_minutes(text: str) -> float:
    words = len(re.findall(r"[A-Za-z0-9']+", text or ""))
    wpm = 150.0
    minutes = words / wpm if wpm else 0.0
    return max(minutes, 0.01)


def _reserve_tts_minutes(profile: PaymentsUserProfile, estimated_minutes: float) -> bool:
    if _is_premium(profile):
        return True
    allowed_before = FREE_AUDIO_MINUTES_LIMIT - float(estimated_minutes)
    if allowed_before < 0:
        return False
    updated = PaymentsUserProfile.objects.filter(
        pk=profile.pk,
        audio_minutes_used__lte=allowed_before
    ).update(audio_minutes_used=F("audio_minutes_used") + float(estimated_minutes))
    return updated == 1


def _refund_tts_minutes(profile: PaymentsUserProfile, estimated_minutes: float) -> None:
    PaymentsUserProfile.objects.filter(pk=profile.pk).update(
        audio_minutes_used=Greatest(F("audio_minutes_used") - float(estimated_minutes), 0.0)
    )


def _to_bool(value, default: bool = True) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return default
    if isinstance(value, (int, float)):
        return value != 0
    normalized = str(value).strip().lower()
    if normalized in {"true", "1", "yes", "y", "on"}:
        return True
    if normalized in {"false", "0", "no", "n", "off", ""}:
        return False
    return default


def _resolve_stored_file_type(file_name: str) -> str:
    ext = file_name.lower().split('.')[-1] if '.' in file_name else ''
    if ext == "pdf":
        return "pdf"
    if ext in {"doc", "docx"}:
        return "docx"
    if ext in {"ppt", "pptx"}:
        return "pptx"
    if ext in {"png", "jpg", "jpeg", "gif", "webp"}:
        return "image"
    if ext in {"mp3", "wav", "ogg", "oga", "m4a"}:
        return "audio"
    return "other"


def _persist_source_upload(user, uploaded_file, purpose: str) -> StoredFile:
    stored = StoredFile.objects.create(
        user=user,
        file=uploaded_file,
        file_type=_resolve_stored_file_type(uploaded_file.name),
        purpose=purpose,
        status="completed",
    )
    # Ensure extraction reads from the start after storage save.
    try:
        uploaded_file.seek(0)
    except Exception:
        pass
    return stored


def _is_default_chat_title(title: str) -> bool:
    normalized = str(title or "").strip().lower()
    return normalized in {"", "new chat", "new study session", "untitled chat"}


def _build_feature_chat_title(source_name: str, suffix: str) -> str:
    filename = str(source_name or "").rsplit("/", 1)[-1]
    stem = filename.rsplit(".", 1)[0] if "." in filename else filename
    stem = re.sub(r"[_\-]+", " ", stem)
    stem = re.sub(r"\s+", " ", stem).strip()
    if not stem:
        stem = "Study Material"
    if len(stem) > 70:
        stem = f"{stem[:67].rstrip()}..."
    title = f"{stem} {suffix}".strip()
    return title[:255]


def _auto_title_feature_chat(chat: Chat | None, source_name: str, suffix: str) -> str | None:
    if not chat:
        return None
    if not _is_default_chat_title(chat.title):
        return chat.title
    new_title = _build_feature_chat_title(source_name, suffix)
    chat.title = new_title
    chat.save()
    return chat.title


def _is_smalltalk_prompt(text: str) -> bool:
    normalized = re.sub(r"\s+", " ", str(text or "").strip().lower())
    if not normalized:
        return True
    if len(normalized) <= 40 and re.fullmatch(r"(hi|hey|hello|yo|sup|thanks|thank you)[!. ]*", normalized):
        return True
    smalltalk_patterns = [
        r"\b(good morning|good afternoon|good evening)\b",
        r"\b(how are you|how's it going|what's up)\b",
        r"\b(can you help|help me|assist me)\b",
        r"\b(hello there|hey there)\b",
    ]
    return any(re.search(pattern, normalized) for pattern in smalltalk_patterns)


def _is_substantive_prompt(text: str) -> bool:
    normalized = re.sub(r"\s+", " ", str(text or "").strip())
    if not normalized:
        return False
    words = re.findall(r"[A-Za-z0-9']+", normalized)
    if len(words) < 5 or len(normalized) < 24:
        return False
    return not _is_smalltalk_prompt(normalized)


def _looks_provisional_chat_title(title: str) -> bool:
    normalized = str(title or "").strip().lower()
    if _is_default_chat_title(normalized):
        return True
    return any(token in normalized for token in ["greeting", "assist", "hello", "new chat"])


def _derive_title_from_prompt(prompt_text: str) -> str:
    normalized = re.sub(r"\s+", " ", str(prompt_text or "").strip())
    if not normalized:
        return "New Chat"
    cleaned = re.sub(r"^[#*\-\d.\s:]+", "", normalized).strip()
    if not cleaned:
        return "New Chat"
    words = cleaned.split(" ")
    snippet = " ".join(words[:7]).strip()
    if not snippet:
        return "New Chat"
    snippet = snippet[0].upper() + snippet[1:] if len(snippet) > 1 else snippet.upper()
    if len(snippet) > 70:
        snippet = f"{snippet[:67].rstrip()}..."
    return snippet


def _maybe_generate_chat_title(chat: Chat | None, prompt_text: str, response_text: str) -> str | None:
    if not chat:
        return None

    prompt_text = str(prompt_text or "").strip()
    response_text = str(response_text or "").strip()
    current_title = str(chat.title or "").strip()
    if not (_is_default_chat_title(current_title) or _looks_provisional_chat_title(current_title)):
        return chat.title
    if not _is_substantive_prompt(prompt_text):
        return chat.title

    conversation = []
    if prompt_text:
        conversation.append({"role": "user", "content": prompt_text})
    if response_text:
        conversation.append({"role": "assistant", "content": response_text})
    if not conversation:
        return chat.title

    title = generate_chat_title_from_openai(conversation)
    if not title:
        fallback = _derive_title_from_prompt(prompt_text)
        if _is_default_chat_title(chat.title) and fallback != "New Chat":
            chat.title = fallback[:255]
            chat.save(update_fields=["title", "updated_at"])
        return chat.title

    cleaned_title = re.sub(r"\s+", " ", str(title).strip().strip("\"'"))
    cleaned_title = re.sub(r"^[#*\-\d.\s:]+", "", cleaned_title).strip()
    if not cleaned_title:
        fallback = _derive_title_from_prompt(prompt_text)
        if _is_default_chat_title(chat.title) and fallback != "New Chat":
            chat.title = fallback[:255]
            chat.save(update_fields=["title", "updated_at"])
        return chat.title

    chat.title = cleaned_title[:255]
    chat.save(update_fields=["title", "updated_at"])
    return chat.title

class AddCollaboratorAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request, chat_id):
        user_ids = request.data.get("user_ids", [])
        access_level = request.data.get("access_level", "view")

        if not isinstance(user_ids, list) or not user_ids:
            return Response({"error": "user_ids must be a non-empty list."}, status=400)

        # Get chat object first
        chat = get_object_or_404(Chat, id=chat_id)

        # ✅ Check ownership before anything else
        if chat.user != request.user:
            return Response({"error": "Only the owner can add collaborators."}, status=403)

        added = []

        for user_id in user_ids:
            try:
                user = User.objects.get(id=user_id)
                if user == request.user:
                    continue  # Skip adding self

                existing = ChatCollaborator.objects.filter(chat=chat, collaborator=user).first()

                if existing:
                    if existing.access_level != access_level:
                        existing.access_level = access_level
                        existing.save()
                    continue  # Already added
                else:
                    ChatCollaborator.objects.create(
                        chat=chat,
                        collaborator=user,
                        access_level=access_level,
                        added_by=request.user,
                        is_approved=False
                    )
                    added.append({"id": user.id, "username": user.username})

            except User.DoesNotExist:
                continue  # Ignore invalid IDs

        return Response({
            "message": f"{len(added)} collaborators added.",
            "collaborators": added
        }, status=200)

class RemoveCollaboratorAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request, chat_id):
        username = request.data.get("username")
        if not username:
            return Response({"error": "Username is required."}, status=400)

        try:
            chat = Chat.objects.get(id=chat_id)
            collaborator = User.objects.get(username=username)
            collab = ChatCollaborator.objects.get(chat=chat, collaborator=collaborator)
        except (Chat.DoesNotExist, User.DoesNotExist, ChatCollaborator.DoesNotExist):
            return Response({"error": "Collaboration not found."}, status=404)

        if request.user != chat.user and request.user != collab.added_by:
            return Response({"error": "You are not authorized to remove this collaborator."}, status=403)

        collab.delete()
        return Response({"message": f"{username} removed from collaborators."}, status=200)

class ChatMessagesAPIView(RetrieveAPIView):
    permission_classes = [IsAuthenticated]

    def get(self, request, chat_id):
        try:
            chat = Chat.objects.get(id=chat_id, user=request.user)
        except Chat.DoesNotExist:
            return Response({"error": "Chat not found."}, status=status.HTTP_404_NOT_FOUND)

        messages = ChatHistory.objects.filter(chat=chat).order_by('created_at')
        serializer = ChatMessageSerializer(messages, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)


class ChatListAPIView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        try:
            owned_chats = Chat.objects.filter(
                user=request.user, is_deleted=False
            )
            shared_chats = Chat.objects.filter(
                collaborators__collaborator=request.user,
                collaborators__is_approved=True,
                is_deleted=False
            )

            chats = (owned_chats | shared_chats).distinct().order_by('-updated_at')
            serializer = ChatListSerializer(chats, many=True)
            return Response(serializer.data, status=200)

        except Exception as e:
            logger.error(f"Error fetching chat list: {str(e)}", exc_info=True)
            return Response({"error": "Unable to fetch chat list."}, status=500)




@method_decorator(csrf_exempt, name='dispatch')
class ChatAPIView(APIView):
    parser_classes = [MultiPartParser, FormParser, JSONParser]

    def post(self, request):
        user = request.user if request.user.is_authenticated else None
        raw_prompt = (request.data.get("prompt") or "").strip()
        prompt = raw_prompt
        chat_id = request.data.get("chat_id")
        guest_id = request.data.get("guest_id")
        study_methods = _parse_study_methods(request.data.get("study_methods"))
        study_custom_prompt = _normalize_study_custom_prompt(request.data.get("study_custom_prompt"))
        ip = get_client_ip(request)

        uploaded_files = request.FILES.getlist("files") if hasattr(request, "FILES") else []
        uploaded_file_descriptors = []

        file_context_chunks = []
        unsupported_files = []
        unreadable_files = []
        max_file_context_chars = 3000
        remaining_chars = max_file_context_chars
        allowed_file_exts = {"pdf", "txt", "doc", "docx", "md", "ppt", "pptx"}

        for uploaded_file in uploaded_files:
            file_ext = uploaded_file.name.lower().split(".")[-1] if "." in uploaded_file.name else ""
            uploaded_file_descriptors.append({
                "name": uploaded_file.name,
                "size": int(getattr(uploaded_file, "size", 0) or 0),
                "ext": file_ext,
                "content_type": getattr(uploaded_file, "content_type", "") or "",
            })
            if remaining_chars <= 0:
                break
            if file_ext not in allowed_file_exts:
                unsupported_files.append(uploaded_file.name)
                continue
            extracted = extract_text_from_file(uploaded_file)
            if not extracted:
                unreadable_files.append(uploaded_file.name)
                continue
            excerpt = extracted[:remaining_chars].strip()
            if not excerpt:
                unreadable_files.append(uploaded_file.name)
                continue
            file_context_chunks.append(f"[{uploaded_file.name}]\n{excerpt}")
            remaining_chars -= len(excerpt)

        if file_context_chunks:
            combined_file_context = "\n\n".join(file_context_chunks)
            prompt = (
                f"{prompt}\n\nReference material from uploaded files:\n{combined_file_context}"
                if prompt
                else f"Use this uploaded study material to help me:\n{combined_file_context}"
            )

        if uploaded_files and not file_context_chunks and not prompt:
            details = []
            if unsupported_files:
                details.append(
                    f"Unsupported file type: {', '.join(unsupported_files)}. Allowed: {', '.join(sorted(allowed_file_exts))}."
                )
            if unreadable_files:
                details.append(f"Could not extract readable text from: {', '.join(unreadable_files)}.")
            if not details:
                details.append("No usable text was found in the uploaded files.")
            return Response({"error": " ".join(details)}, status=status.HTTP_400_BAD_REQUEST)

        logger.debug(f"[ChatAPI] Request received - user: {user}, prompt length: {len(prompt) if prompt else 0}, chat_id: {chat_id}, guest_id: {guest_id}")
        logger.debug(f"[ChatAPI] Request data: {dict(request.data)}")

        if not prompt:
            logger.warning(f"[ChatAPI] Missing prompt - request.data: {dict(request.data)}")
            return Response(
                {"error": "Please enter a message to begin."},
                status=status.HTTP_400_BAD_REQUEST
            )

        # Validate prompt length (prevent DoS and cost abuse)
        MAX_PROMPT_LENGTH = 4000
        if len(prompt) > MAX_PROMPT_LENGTH:
            return Response(
                {"error": f"Message too long. Maximum {MAX_PROMPT_LENGTH} characters allowed."},
                status=status.HTTP_400_BAD_REQUEST
            )

        # Guest user flow
        if not user:
            if not guest_id:
                guest_uuid = uuid4()
            else:
                try:
                    guest_uuid = UUID(guest_id)
                except ValueError:
                    guest_uuid = uuid4()

            guest_tracker, _ = GuestChatTracker.objects.get_or_create(guest_id=guest_uuid)
            ip_tracker, _ = GuestIPTracker.objects.get_or_create(ip_address=ip)

            if guest_tracker.count >= 10 or ip_tracker.count >= 10:
                return Response({
                    "limit_exceeded": True,
                    "message": "You've reached your guest chat limit. Please log in to continue."
                }, status=status.HTTP_200_OK)

        # Get or create chat
        try:
            if user:
                if chat_id:
                    chat = Chat.objects.get(id=chat_id)
                    
                    # Ownership or collaborator check
                    if chat.user != user:
                        collaborator = ChatCollaborator.objects.filter(chat=chat, collaborator=user, is_approved=True).first()
                        if not collaborator:
                            return Response({"error": "You do not have access to this chat."}, status=403)
                        if collaborator.access_level != "edit":
                            return Response({"error": "You don't have permission to modify this chat."}, status=403)
                else:
                    chat = Chat.objects.create(user=user, title="New Chat")
            else:
                chat = None  # Guests don't persist chat
        except Chat.DoesNotExist:
            return Response({
                "error": "We couldn't find that chat session."
            }, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            logger.error(f"[Chat Init] {e}", exc_info=True)
            return Response({
                "error": "We had trouble starting your chat. Please try again."
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        system_prompt = _build_chat_system_prompt(study_methods, study_custom_prompt)
        
        # Build conversation history
        messages = [{"role": "system", "content": system_prompt}]

        # Add conversation history if chat exists and user is authenticated
        if user and chat:
            recent_messages = ChatHistory.objects.filter(chat=chat).order_by('-created_at')[:20]  # Last 20 messages for context
            # Reverse to get chronological order and truncate long messages
            for msg in reversed(recent_messages):
                # Truncate individual messages to preserve token limits (max 500 chars per message)
                prompt_content = msg.prompt[:500] if len(msg.prompt) > 500 else msg.prompt
                response_content = msg.response[:500] if len(msg.response) > 500 else msg.response
                messages.append({"role": "user", "content": prompt_content})
                messages.append({"role": "assistant", "content": response_content})

        # Add current user prompt
        messages.append({"role": "user", "content": prompt})

        # Get response style settings from request
        response_style = request.data.get("response_style", "balanced")

        # Configure OpenAI parameters based on response style
        style_configs = {
            "concise": {
                "temperature": 0.3,
                "max_tokens": 600,
            },
            "balanced": {
                "temperature": 0.5,
                "max_tokens": 1400,
            },
            "detailed": {
                "temperature": 0.7,
                "max_tokens": 2600,
            },
        }

        config = style_configs.get(response_style, style_configs["balanced"])

        try:
            client = get_openai_client()
            reply_text = _generate_chat_completion_text(
                client=client,
                model=getattr(settings, "CHAT_OPENAI_MODEL", "gpt-4o-mini"),
                messages=messages,
                temperature=config["temperature"],
                max_tokens=config["max_tokens"],
            )
            bot_reply = _sanitize_response_text(reply_text)
        except Exception as e:
            logger.error(f"[OpenAI Error] {e}", exc_info=True)
            return Response({
                "error": "We're having trouble responding. Please try again shortly."
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        # Save chat history
        chat_title = None
        if user:
            try:
                prompt_type = "attachment" if uploaded_files else "text"
                prompt_metadata = {
                    "caption": raw_prompt,
                    "files": uploaded_file_descriptors,
                    "study_methods": study_methods,
                    "study_custom_prompt": study_custom_prompt,
                } if uploaded_files else {
                    "study_methods": study_methods,
                    "study_custom_prompt": study_custom_prompt,
                }
                ChatHistory.objects.create(
                    chat=chat,
                    user=user,
                    prompt=raw_prompt,
                    response=bot_reply,
                    prompt_type=prompt_type,
                    prompt_metadata=prompt_metadata,
                    response_type="text",
                    context="file" if uploaded_files else "general"
                )
                chat_title = _maybe_generate_chat_title(chat, raw_prompt, bot_reply)
            except Exception as e:
                logger.warning(f"[History Save Fail] {e}", exc_info=True)

        # Update guest usage
        if not user:
            guest_tracker.count += 1
            guest_tracker.save()
            ip_tracker.count += 1
            ip_tracker.save()

        response_payload = {
            "chat_id": chat.id if chat else None,
            "response": bot_reply,
            "chat_title": chat_title,
            "limit_exceeded": False,
            "uploaded_files": uploaded_file_descriptors,
        }
        if not user:
            response_payload["guest_id"] = str(guest_uuid)

        return Response(response_payload, status=status.HTTP_200_OK)


class ChatStreamAPIView(APIView):
    """Streaming version of ChatAPIView for reduced perceived latency"""
    parser_classes = [MultiPartParser, FormParser, JSONParser]

    def post(self, request):
        user = request.user if request.user.is_authenticated else None
        raw_prompt = (request.data.get("prompt") or "").strip()
        prompt = raw_prompt
        chat_id = request.data.get("chat_id")
        guest_id = request.data.get("guest_id")
        study_methods = _parse_study_methods(request.data.get("study_methods"))
        study_custom_prompt = _normalize_study_custom_prompt(request.data.get("study_custom_prompt"))
        ip = get_client_ip(request)

        uploaded_files = request.FILES.getlist("files") if hasattr(request, "FILES") else []
        
        # For streaming, we don't extract file context - keep it simple for now
        # File processing can be added later if needed

        if not prompt:
            return Response(
                {"error": "Please enter a message to begin."},
                status=status.HTTP_400_BAD_REQUEST
            )

        MAX_PROMPT_LENGTH = 4000
        if len(prompt) > MAX_PROMPT_LENGTH:
            return Response(
                {"error": f"Message too long. Maximum {MAX_PROMPT_LENGTH} characters allowed."},
                status=status.HTTP_400_BAD_REQUEST
            )

        # Guest user flow
        if not user:
            if not guest_id:
                guest_uuid = uuid4()
            else:
                try:
                    guest_uuid = UUID(guest_id)
                except ValueError:
                    guest_uuid = uuid4()

            guest_tracker, _ = GuestChatTracker.objects.get_or_create(guest_id=guest_uuid)
            ip_tracker, _ = GuestIPTracker.objects.get_or_create(ip_address=ip)

            if guest_tracker.count >= 10 or ip_tracker.count >= 10:
                return Response({
                    "limit_exceeded": True,
                    "message": "You've reached your guest chat limit. Please log in to continue."
                }, status=status.HTTP_200_OK)

        # Get or create chat
        try:
            if user:
                if chat_id:
                    chat = Chat.objects.get(id=chat_id)
                    if chat.user != user:
                        collaborator = ChatCollaborator.objects.filter(
                            chat=chat, collaborator=user, is_approved=True
                        ).first()
                        if not collaborator:
                            return Response({"error": "You do not have access to this chat."}, status=403)
                        if collaborator.access_level != "edit":
                            return Response({"error": "You don't have permission to modify this chat."}, status=403)
                else:
                    chat = Chat.objects.create(user=user, title="New Chat")
            else:
                chat = None
        except Chat.DoesNotExist:
            return Response({"error": "We couldn't find that chat session."}, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            logger.error(f"[Chat Stream Init] {e}", exc_info=True)
            return Response({"error": "We had trouble starting your chat. Please try again."}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        system_prompt = _build_chat_system_prompt(study_methods, study_custom_prompt)

        # Build conversation messages
        messages = [{"role": "system", "content": system_prompt}]

        # Add conversation history (reduced for streaming performance)
        if user and chat:
            recent_messages = ChatHistory.objects.filter(chat=chat).order_by('-created_at')[:10]
            for msg in reversed(recent_messages):
                prompt_content = msg.prompt[:500] if len(msg.prompt) > 500 else msg.prompt
                response_content = msg.response[:500] if len(msg.response) > 500 else msg.response
                messages.append({"role": "user", "content": prompt_content})
                messages.append({"role": "assistant", "content": response_content})

        messages.append({"role": "user", "content": prompt})

        # Get response style settings
        response_style = request.data.get("response_style", "balanced")
        style_configs = {
            "concise": {"temperature": 0.3, "max_tokens": 600},
            "balanced": {"temperature": 0.5, "max_tokens": 1400},
            "detailed": {"temperature": 0.7, "max_tokens": 2600},
        }
        config = style_configs.get(response_style, style_configs["balanced"])

        def generate():
            """Generator function for streaming response"""
            try:
                client = get_openai_client()

                meta_payload = {
                    "type": "meta",
                    "chat_id": chat.id if chat else None,
                    "guest_id": str(guest_uuid) if not user else None,
                    "uploaded_files": [],
                }
                yield _sse_payload(meta_payload)

                # Stream the response
                stream = client.chat.completions.create(
                    model=getattr(settings, "CHAT_OPENAI_MODEL", "gpt-4o-mini"),
                    messages=messages,
                    temperature=config["temperature"],
                    max_tokens=config["max_tokens"],
                    stream=True,
                )
                
                full_response = ""
                for chunk in stream:
                    if chunk.choices and chunk.choices[0].delta.content:
                        content = chunk.choices[0].delta.content
                        full_response += content
                        yield _sse_payload({"type": "delta", "delta": content})

                bot_reply = _sanitize_response_text(full_response)
                chat_title = None
                if user:
                    try:
                        ChatHistory.objects.create(
                            chat=chat,
                            user=user,
                            prompt=raw_prompt,
                            response=bot_reply,
                            prompt_type="text",
                            prompt_metadata={
                                "study_methods": study_methods,
                                "study_custom_prompt": study_custom_prompt,
                            },
                            response_type="text",
                            context="general",
                        )
                        chat_title = _maybe_generate_chat_title(chat, raw_prompt, bot_reply)
                    except Exception as e:
                        logger.warning(f"[Chat Stream History Save Fail] {e}", exc_info=True)
                else:
                    guest_tracker.count += 1
                    guest_tracker.save(update_fields=["count"])
                    ip_tracker.count += 1
                    ip_tracker.save(update_fields=["count"])

                yield _sse_payload({"type": "final", "response": bot_reply})
                yield _sse_payload({
                    "type": "done",
                    "chat_id": chat.id if chat else None,
                    "response": bot_reply,
                    "chat_title": chat_title,
                    "limit_exceeded": False,
                    "guest_id": str(guest_uuid) if not user else None,
                    "uploaded_files": [],
                })
            except Exception as e:
                logger.error(f"[Chat Stream Error] {e}", exc_info=True)
                yield _sse_payload({
                    "type": "error",
                    "error": "We're having trouble responding. Please try again shortly."
                })

        response = StreamingHttpResponse(generate(), content_type="text/event-stream")
        response["Cache-Control"] = "no-cache"
        response["X-Accel-Buffering"] = "no"
        return response


class ChatDetailView(APIView):
    permission_classes = [IsAuthenticated]

    def patch(self, request, chat_id):
        new_title = request.data.get("title")
        if not new_title:
            return Response({"error": "Title is required."}, status=400)

        try:
            chat = Chat.objects.get(id=chat_id, user=request.user)
        except Chat.DoesNotExist:
            return Response({"error": "Chat not found."}, status=404)

        chat.title = new_title
        chat.save()
        return Response({"message": "Chat title updated successfully.", "title": new_title})

    def delete(self, request, chat_id):
        try:
            chat = Chat.objects.get(id=chat_id, user=request.user)
        except Chat.DoesNotExist:
            return Response({"error": "Chat not found."}, status=404)

        chat.is_deleted = True
        chat.save()
        return Response({"message": "Chat deleted successfully."}, status=204)





class CommitChatView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        chat_id = request.data.get("chat_id")

        if not chat_id:
            return Response({"error": "Chat ID is required."}, status=400)

        try:
            chat = Chat.objects.get(id=chat_id)
        except Chat.DoesNotExist:
            return Response({"error": "Chat not found."}, status=404)

        user = request.user
        is_owner = chat.user == user
        has_edit_access = ChatCollaborator.objects.filter(
            chat=chat,
            collaborator=user,
            is_approved=True,
            access_level="edit"
        ).exists()

        if not is_owner and not has_edit_access:
            return Response({"error": "You do not have permission to modify this chat."}, status=403)

        history = list(ChatHistory.objects.filter(chat=chat).order_by("created_at")[:6])
        if not history:
            return Response(
                {"error": "At least one message is required to generate a title."},
                status=400
            )

        seed_pair = None
        has_substantive_prompt = False
        for item in history:
            prompt = str(item.prompt or "").strip()
            if not prompt:
                continue
            if _is_substantive_prompt(prompt):
                has_substantive_prompt = True
                seed_pair = item
                break
            if seed_pair is None:
                seed_pair = item

        if seed_pair is None:
            seed_pair = history[0]

        conversation = []
        prompt_text = str(seed_pair.prompt or "").strip()
        response_text = str(seed_pair.response or "").strip()
        if prompt_text:
            conversation.append({"role": "user", "content": prompt_text})
        if response_text:
            conversation.append({"role": "assistant", "content": response_text})

        if not conversation:
            return Response(
                {"error": "No usable message content found to generate a title."},
                status=400
            )

        generated_title = generate_chat_title_from_openai(conversation)
        title = generated_title or _derive_title_from_prompt(prompt_text)

        first_prompt = str(history[0].prompt or "").strip()
        first_prompt_snippet = (
            f"{first_prompt[:60].rstrip()}..." if len(first_prompt) > 60 else first_prompt
        )
        current_title = str(chat.title or "").strip()
        history_count = ChatHistory.objects.filter(chat=chat).count()
        should_replace = (
            _is_default_chat_title(current_title)
            or (first_prompt_snippet and current_title.lower() == first_prompt_snippet.lower())
            or (_looks_provisional_chat_title(current_title) and has_substantive_prompt and history_count <= 8)
        )

        if should_replace:
            chat.title = title
            chat.save(update_fields=["title", "updated_at"])

        return Response({"chat_id": chat.id, "title": chat.title}, status=200)


from rest_framework.views import APIView

class UpdateChatTitleAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request, chat_id):
        chat = get_object_or_404(Chat, id=chat_id)

        # Owner check
        if chat.user == request.user:
            allowed = True
        else:
            collab = ChatCollaborator.objects.filter(
                chat=chat,
                collaborator=request.user,
                is_approved=True,
                access_level="edit"  # <-- must be edit
            ).first()
            allowed = bool(collab)

        if not allowed:
            return Response({"error": "Permission denied."}, status=403)

        title = request.data.get("title")
        if title:
            chat.title = title
            chat.save()

        return Response({"message": "Chat updated."})

class DeleteChatAPIView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, chat_id):
        user = request.user

        try:
            chat = Chat.objects.get(id=chat_id)
        except Chat.DoesNotExist:
            return Response({"error": "Chat not found."}, status=404)

        # Check ownership or edit-level collaborator
        is_owner = chat.user == user
        has_edit_access = ChatCollaborator.objects.filter(
            chat=chat,
            collaborator=user,
            is_approved=True,
            access_level="edit"
        ).exists()

        if not is_owner and not has_edit_access:
            return Response({"error": "You do not have permission to delete this chat."}, status=403)

        chat.is_deleted = True
        chat.save()
        return Response({"message": "Chat deleted successfully."}, status=200)


class StartNewChatAPIView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        chat = Chat.objects.create(user=request.user, title="New Chat")
        return Response({"chat_id": chat.id, "message": "New chat started."}, status=201)


from django.db.models import Q

class ListUsersAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def get(self, request):
        query = request.GET.get("q", "").strip()
        users = User.objects.exclude(id=request.user.id)
        if query:
            users = users.filter(Q(username__icontains=query) | Q(email__icontains=query))
        data = [{"id": u.id, "username": u.username, "email": u.email} for u in users]
        return Response(data)


class ShareChatAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request, chat_id):
        chat = get_object_or_404(Chat, id=chat_id, user=request.user)
        user_ids = request.data.get("user_ids", [])
        access_level = request.data.get("access_level", "view")

        if access_level not in ["view", "edit"]:
            return Response({"error": "Invalid access level."}, status=status.HTTP_400_BAD_REQUEST)

        added = []
        skipped = []
        for uid in user_ids:
            user = User.objects.filter(id=uid).first()
            if not user:
                skipped.append(uid)
                continue

            obj, created = ChatCollaborator.objects.get_or_create(
                chat=chat, collaborator=user,
                defaults={"added_by": request.user, "access_level": access_level, "is_approved": False}
            )
            if created:
                added.append(user.username)
            else:
                skipped.append(user.username)

        return Response({
            "added": added,
            "skipped": skipped
        })


class ApproveCollaborationAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request, chat_id):
        collab = get_object_or_404(
            ChatCollaborator, chat_id=chat_id, collaborator=request.user, is_approved=False
        )
        collab.is_approved = True
        collab.save()
        return Response({"message": "Collaboration approved."})



class EmailShareAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request, chat_id):
        email = request.data.get("email")
        access_level = request.data.get("access_level", "view")
        user = User.objects.filter(email=email).first()
        chat = get_object_or_404(Chat, id=chat_id, user=request.user)

        if not user:
            return Response({"error": "User with this email not found."}, status=status.HTTP_404_NOT_FOUND)

        if user == request.user:
            return Response({"error": "You cannot share the chat with yourself."}, status=status.HTTP_400_BAD_REQUEST)

        _, created = ChatCollaborator.objects.get_or_create(
            chat=chat, collaborator=user,
            defaults={"added_by": request.user, "access_level": access_level, "is_approved": False}
        )
        # Replace this with actual email logic if needed.
        return Response({"message": f"Chat shared with {user.username}."})




class ListCollaboratorsAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def get(self, request, chat_id):
        chat = get_object_or_404(Chat, id=chat_id)

        # Check access
        is_owner = chat.user == request.user
        is_collaborator = ChatCollaborator.objects.filter(
            chat=chat, collaborator=request.user, is_approved=True
        ).exists()

        if not is_owner and not is_collaborator:
            return Response({"error": "Not authorized."}, status=403)

        # Get actual collaborators
        collaborators = list(ChatCollaborator.objects.filter(chat=chat))

        # Add the owner as a pseudo-collaborator
        owner_user = chat.user
        owner_data = {
            "id": -1,  # Indicates synthetic collaborator
            "chat": chat.id,
            "collaborator": {
                "id": owner_user.id,
                "username": owner_user.username,
                "email": owner_user.email
            },
            "added_by": {
                "id": owner_user.id,
                "username": owner_user.username,
                "email": owner_user.email
            },
            "access_level": "edit",
            "is_approved": True,
            "added_at": chat.created_at,
            "is_owner": True
        }

        # Serialize real collaborators
        serializer = ChatCollaboratorSerializer(collaborators, many=True)
        data = serializer.data

        # Add synthetic owner entry
        data.insert(0, owner_data)

        return Response(data, status=200)


class GuestChatStatusAPIView(APIView):
    permission_classes = [permissions.AllowAny]

    def get(self, request):
        guest_id = request.query_params.get("guest_id")
        ip = get_client_ip(request)

        guest_count = 0
        ip_count = 0

        if guest_id:
            try:
                guest_uuid = UUID(str(guest_id))
                guest_tracker = GuestChatTracker.objects.filter(guest_id=guest_uuid).first()
                guest_count = guest_tracker.count if guest_tracker else 0
            except (ValueError, TypeError):
                guest_count = 0

        ip_tracker = GuestIPTracker.objects.filter(ip_address=ip).first()
        ip_count = ip_tracker.count if ip_tracker else 0

        limit = 10
        limit_exceeded = guest_count >= limit or ip_count >= limit
        remaining = max(0, limit - max(guest_count, ip_count))

        return Response(
            {
                "limit_exceeded": limit_exceeded,
                "remaining": remaining,
                "message": "You've reached your guest chat limit. Please log in to continue."
                if limit_exceeded
                else "",
            },
            status=200,
        )


class PendingCollaborationsAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def get(self, request):
        collaborations = ChatCollaborator.objects.filter(
            collaborator=request.user,
            is_approved=False
        ).select_related("chat", "chat__user")  # important

        data = [{
            "id": c.chat.id,
            "title": c.chat.title,
            "owner": {
                "id": c.chat.user.id,
                "username": c.chat.user.username,
                "email": c.chat.user.email
            }
        } for c in collaborations]

        return Response(data)

class RejectCollaborationAPIView(APIView):
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request, chat_id):
        chat = get_object_or_404(Chat, id=chat_id)

        try:
            collab = ChatCollaborator.objects.get(chat=chat, collaborator=request.user, is_approved=False)
        except ChatCollaborator.DoesNotExist:
            return Response({"error": "No pending invitation found."}, status=404)

        collab.delete()
        return Response({"success": "Collaboration rejected."})


# views.py - Add these new views


class TextToAudioView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = [MultiPartParser, FormParser, JSONParser]
    
    def post(self, request):
        text = request.data.get('text')
        if not text and 'file' in request.FILES:
            upload = request.FILES['file']
            allowed_types = {'pdf', 'txt', 'doc', 'docx', 'md'}
            ext = upload.name.lower().split('.')[-1] if '.' in upload.name else ''
            if ext not in allowed_types:
                return Response(
                    {"error": f"Invalid file type. Allowed: {', '.join(sorted(allowed_types))}"},
                    status=400
                )
            max_file_bytes = max(1024, int(getattr(settings, "TTS_SOURCE_MAX_FILE_BYTES", 10 * 1024 * 1024)))
            if upload.size > max_file_bytes:
                return Response(
                    {"error": f"File too large. Maximum size is {max_file_bytes // (1024 * 1024)}MB."},
                    status=400
                )
            text = extract_text_from_file(upload)
        if isinstance(text, str):
            text = text.strip()
        if not text:
            return Response({"error": "No text provided"}, status=400)
        chat_id = request.data.get("chat_id")
        voice = request.data.get("voice", "alloy")
        allowed_voices = {"alloy", "echo", "fable", "onyx", "nova", "shimmer"}
        if voice not in allowed_voices:
            return Response(
                {"error": f"Invalid voice. Allowed: {', '.join(sorted(allowed_voices))}"},
                status=400
            )
        
        profile = _get_payments_profile(request.user)
        premium_active = _is_premium(profile)
        max_tts_length = int(
            getattr(
                settings,
                "TTS_MAX_TEXT_LENGTH_PRO" if premium_active else "TTS_MAX_TEXT_LENGTH_FREE",
                5000 if premium_active else 2000,
            )
        )
        max_tts_length = max(100, max_tts_length)
        if len(text) > max_tts_length:
            return Response(
                {"error": f"Text too long. Maximum {max_tts_length} characters allowed for text-to-speech."},
                status=400
            )

        estimated_minutes = _estimate_tts_minutes(text)
        quota_reserved = _reserve_tts_minutes(profile, estimated_minutes)
        if not quota_reserved:
            return Response(
                {"error": "Free audio limit reached. Upgrade to premium."},
                status=429
            )
        
        try:
            started_at = time.monotonic()
            tts_model = getattr(settings, "TTS_OPENAI_MODEL", "gpt-4o-mini-tts")
            chat = _get_chat_for_write(request.user, chat_id)
            # Generate speech using OpenAI TTS
            client = get_openai_client()
            response = client.audio.speech.create(
                model=tts_model,
                voice=voice,
                input=text
            )
            openai_elapsed = time.monotonic() - started_at

            audio_bytes = response.read() if hasattr(response, "read") else response.content
            tts_obj = None
            audio_url = None
            should_persist_audio = bool(TTS_PERSIST_AUDIO_FILES or (chat and TTS_PERSIST_CHAT_HISTORY))
            if should_persist_audio:
                tts_obj = TextToSpeech.objects.create(
                    user=request.user,
                    text=text
                )
                # `audio_file` has upload_to='text_to_speech/', so save only the filename.
                file_name = f"{tts_obj.id}.mp3"
                tts_obj.audio_file.save(file_name, ContentFile(audio_bytes))
                audio_url = request.build_absolute_uri(tts_obj.audio_file.url)
            elif TTS_RETURN_INLINE_AUDIO:
                audio_url = f"data:audio/mpeg;base64,{base64.b64encode(audio_bytes).decode('ascii')}"
            else:
                # Ensure clients can still play in environments where persistence is disabled.
                audio_url = f"data:audio/mpeg;base64,{base64.b64encode(audio_bytes).decode('ascii')}"

            total_elapsed = time.monotonic() - started_at
            logger.info(
                f"[TTS] model={tts_model} chars={len(text)} openai_s={openai_elapsed:.2f} total_s={total_elapsed:.2f}"
            )

            chat_title = None
            if chat and TTS_PERSIST_CHAT_HISTORY:
                ChatHistory.objects.create(
                    chat=chat,
                    user=request.user,
                    prompt="",
                    response="I've generated audio for your text. Click play to listen.",
                    prompt_type="text",
                    response_type="audio",
                    response_metadata={
                        "tts_id": tts_obj.id if tts_obj else None,
                        "title": "Generated Audio",
                        "audio_url": audio_url,
                        "voice": voice,
                        "transcript": text,
                    },
                    context="audio",
                )
                source_label = request.FILES['file'].name if 'file' in request.FILES else (text[:60] if text else "Text to Speech")
                chat_title = _auto_title_feature_chat(chat, source_label, "Audio")
            elif chat:
                # Keep title freshness in turbo mode without waiting on ChatHistory writes.
                source_label = request.FILES['file'].name if 'file' in request.FILES else (text[:60] if text else "Text to Speech")
                chat_title = _auto_title_feature_chat(chat, source_label, "Audio")

            return Response({
                "id": tts_obj.id if tts_obj else None,
                "audio_url": audio_url,
                "chat_title": chat_title,
                "persisted": bool(tts_obj),
            })
            
        except Exception as e:
            if quota_reserved and not premium_active:
                _refund_tts_minutes(profile, estimated_minutes)
            logger.error(f"Text-to-speech failed: {str(e)}")
            return Response({"error": "Text-to-speech conversion failed"}, status=500)


class AudioFileByIdView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request, tts_id: int):
        tts_obj = TextToSpeech.objects.filter(id=tts_id).first()
        if not tts_obj or not tts_obj.audio_file:
            return Response({"error": "Audio file not found."}, status=404)
        if tts_obj.user_id != request.user.id:
            return Response({"error": "You do not have access to this audio file."}, status=403)

        try:
            audio_stream = tts_obj.audio_file.open("rb")
        except FileNotFoundError:
            return Response({"error": "Audio file not found."}, status=404)
        except Exception:
            return Response({"error": "Failed to open audio file."}, status=500)

        filename = tts_obj.audio_file.name.rsplit("/", 1)[-1] or f"{tts_id}.mp3"
        return FileResponse(audio_stream, content_type="audio/mpeg", as_attachment=False, filename=filename)


class AudioFileDownloadView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request, tts_id: int):
        tts_obj = TextToSpeech.objects.filter(id=tts_id).first()
        if not tts_obj or not tts_obj.audio_file:
            return Response({"error": "Audio file not found."}, status=404)
        if tts_obj.user_id != request.user.id:
            return Response({"error": "You do not have access to this audio file."}, status=403)

        profile = _get_payments_profile(request.user)
        if not _is_premium(profile):
            return Response({"error": "Downloading audio is available on the Pro plan."}, status=403)

        try:
            audio_stream = tts_obj.audio_file.open("rb")
        except FileNotFoundError:
            return Response({"error": "Audio file not found."}, status=404)
        except Exception:
            return Response({"error": "Failed to open audio file."}, status=500)

        filename = tts_obj.audio_file.name.rsplit("/", 1)[-1] or f"{tts_id}.mp3"
        return FileResponse(audio_stream, content_type="audio/mpeg", as_attachment=True, filename=filename)


import io
import PyPDF2
import time
from django.core.files.uploadedfile import UploadedFile


def extract_text_from_file(file_obj) -> str:
    """Extract text from uploaded file (PDF, TXT, DOCX, PPTX, etc.)"""
    file_obj.seek(0)
    
    # Get file extension
    name = file_obj.name.lower()
    
    if name.endswith('.pdf'):
        try:
            pdf_reader = PyPDF2.PdfReader(file_obj)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return ""
    
    elif name.endswith('.txt') or name.endswith('.md'):
        return file_obj.read().decode('utf-8', errors='ignore')
    
    elif name.endswith('.docx'):
        try:
            from docx import Document
            doc = Document(file_obj)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return ""
    
    elif name.endswith('.pptx'):
        try:
            from pptx import Presentation
            prs = Presentation(file_obj)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return ""
    
    else:
        # Try to read as plain text
        try:
            content = file_obj.read()
            if isinstance(content, bytes):
                return content.decode('utf-8', errors='ignore')
            return str(content)
        except Exception as e:
            logger.error(f"File reading failed: {e}")
            return ""


def _prepare_llm_context(text: str, max_chars: int) -> str:
    """Reduce token load while preserving beginning/middle/end coverage."""
    normalized = re.sub(r"\s+", " ", (text or "")).strip()
    if len(normalized) <= max_chars:
        return normalized

    head_len = int(max_chars * 0.5)
    middle_len = int(max_chars * 0.3)
    tail_len = max_chars - head_len - middle_len

    middle_start = max((len(normalized) - middle_len) // 2, 0)
    middle = normalized[middle_start:middle_start + middle_len]
    head = normalized[:head_len]
    tail = normalized[-tail_len:] if tail_len > 0 else ""
    return f"{head}\n...\n{middle}\n...\n{tail}"


def _sample_evenly(items: list, limit: int) -> list:
    if limit <= 0 or len(items) <= limit:
        return items
    if limit == 1:
        return [items[0]]

    last_index = len(items) - 1
    selected = []
    seen = set()
    for step in range(limit):
        idx = round((step * last_index) / (limit - 1))
        if idx in seen:
            continue
        seen.add(idx)
        selected.append(items[idx])
    return selected


def _extract_pdf_page_texts(file_obj, max_pages: int | None = None) -> list[tuple[int, str]]:
    file_obj.seek(0)
    try:
        pdf_reader = PyPDF2.PdfReader(file_obj)
        page_entries: list[tuple[int, str]] = []
        for page_index, page in enumerate(pdf_reader.pages):
            # Hard-stop by physical page index to bound extraction latency.
            if max_pages is not None and page_index >= max_pages:
                break
            text = (page.extract_text() or "").strip()
            if text:
                page_entries.append((page_index + 1, text))
        return page_entries
    except Exception as e:
        logger.error(f"PDF page extraction failed: {e}")
        return []


def _prepare_quiz_source_context(file_obj, max_chars: int, max_pages: int) -> str:
    name = str(getattr(file_obj, "name", "") or "").lower()
    if name.endswith(".pdf"):
        page_entries = _extract_pdf_page_texts(file_obj, max_pages=max_pages)
        if not page_entries:
            return ""

        sampled_entries = _sample_evenly(page_entries, min(len(page_entries), 24))
        snippet_budget = max(300, max_chars // max(1, len(sampled_entries)))
        sections = []
        for page_number, text in sampled_entries:
            normalized = re.sub(r"\s+", " ", text).strip()
            if not normalized:
                continue
            sections.append(f"[Page {page_number}]\n{normalized[:snippet_budget].strip()}")

        combined = "\n\n".join(sections).strip()
        return _prepare_llm_context(combined, max_chars)

    text = extract_text_from_file(file_obj)
    return _prepare_llm_context(text, max_chars)


def _sha256_uploaded_file(file_obj) -> str:
    hasher = hashlib.sha256()
    try:
        file_obj.seek(0)
        if hasattr(file_obj, "chunks"):
            for chunk in file_obj.chunks():
                hasher.update(chunk)
        else:
            hasher.update(file_obj.read())
    finally:
        try:
            file_obj.seek(0)
        except Exception:
            pass
    return hasher.hexdigest()


def _normalize_quiz_options(options) -> list[str]:
    if not isinstance(options, list):
        options = []
    cleaned = [str(opt).strip() for opt in options if str(opt).strip()]
    default_options = ["Option A", "Option B", "Option C", "Option D"]
    cleaned = cleaned[:4]
    while len(cleaned) < 4:
        cleaned.append(default_options[len(cleaned)])
    return cleaned


def _normalize_option_text(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", str(text or "").lower()).strip()


def _pick_option_from_hint(options: list[str], hint: str) -> int | None:
    normalized_hint = _normalize_option_text(hint)
    if not normalized_hint:
        return None

    best_idx = None
    best_score = 0
    for idx, opt in enumerate(options):
        normalized_opt = _normalize_option_text(opt)
        if not normalized_opt:
            continue
        if normalized_opt == normalized_hint:
            return idx
        if normalized_opt in normalized_hint:
            score = len(normalized_opt)
            if score > best_score:
                best_score = score
                best_idx = idx
    return best_idx


def _resolve_mcq_correct_index(raw_correct, options: list[str], correct_text: str, explanation: str) -> int:
    option_count = len(options)
    if option_count == 0:
        return 0

    hint_match = _pick_option_from_hint(options, correct_text)
    if hint_match is None:
        hint_match = _pick_option_from_hint(options, explanation)
    if hint_match is not None:
        return hint_match

    if isinstance(raw_correct, str):
        token = raw_correct.strip().upper()
        letter_map = {"A": 0, "B": 1, "C": 2, "D": 3}
        if token in letter_map and letter_map[token] < option_count:
            return letter_map[token]
        if token.isdigit():
            raw_correct = int(token)

    if isinstance(raw_correct, int):
        # Preferred format is 0-based index.
        if 0 <= raw_correct < option_count:
            return raw_correct
        # Tolerate 1-based index from model drift.
        if 1 <= raw_correct <= option_count:
            return raw_correct - 1

    return 0


def _infer_question_concept(question_text: str) -> str:
    text = re.sub(r"[^a-zA-Z0-9\s]", " ", str(question_text or "")).lower()
    tokens = [t for t in text.split() if len(t) > 3]
    if not tokens:
        return "Core Concept"
    return " ".join(tokens[:2]).title()


def _normalize_summary_payload(summary_data: dict, include_key_terms: bool) -> dict:
    if not isinstance(summary_data, dict):
        summary_data = {}

    summary_content = summary_data.get("summary", "")
    if isinstance(summary_content, dict):
        parts = []
        for section, content in summary_content.items():
            if content:
                parts.append(f"## {section}\n{content}")
        summary_content = "\n\n".join(parts)
    summary_content = _sanitize_response_text(str(summary_content or "").strip())

    takeaways = summary_data.get("takeaways", [])
    if not isinstance(takeaways, list):
        takeaways = []
    takeaways = [str(item).strip() for item in takeaways if str(item).strip()][:8]

    key_terms = []
    if include_key_terms:
        raw_terms = summary_data.get("keyTerms", [])
        if isinstance(raw_terms, list):
            for term_item in raw_terms[:12]:
                if not isinstance(term_item, dict):
                    continue
                term = str(term_item.get("term", "")).strip()
                definition = str(term_item.get("definition", "")).strip()
                if term and definition:
                    key_terms.append({"term": term, "definition": definition})

    return {
        "summary": summary_content,
        "takeaways": takeaways,
        "keyTerms": key_terms,
    }


def _quiz_difficulty_instructions(level: str) -> str:
    profiles = {
        "easy": (
            "- Focus on direct recall of key definitions, facts, and simple concepts.\n"
            "- Use short, clear question wording with one obvious best answer.\n"
            "- Avoid multi-step reasoning and avoid tricky distractors."
        ),
        "medium": (
            "- Focus on understanding and application, not simple recall.\n"
            "- Require 1-2 reasoning steps (e.g., interpretation, comparison, choosing best method).\n"
            "- Use plausible distractors that test common misunderstandings."
        ),
        "hard": (
            "- Focus on analysis, synthesis, and edge cases from the content.\n"
            "- Require multi-step reasoning or careful discrimination between close options.\n"
            "- Use strong, close distractors; avoid ambiguity but make questions genuinely challenging."
        ),
    }
    return profiles.get(level, profiles["medium"])


def _quiz_format_instructions(question_type: str, difficulty: str) -> str:
    base = {
        "mcq": "Generate multiple choice questions with exactly 4 options each. Mark the correct answer with the index (0-3).",
        "tf": "Generate true/false questions. Set options to ['True','False'] and correctAnswer to 0 (True) or 1 (False).",
        "fill": "Generate fill-in-the-blank questions with a clear blank. Set options to empty array [] and correctAnswer to 0.",
        "essay": "Generate essay/theory questions requiring written reasoning. Set options to empty array [] and correctAnswer to 0.",
    }
    difficulty_by_type = {
        "mcq": {
            "easy": "- Keep stems short and straightforward.\n- Avoid negative wording.\n- Distractors should be clearly wrong for prepared students.",
            "medium": "- Mix concept understanding with basic application.\n- Distractors should be plausible.\n- Avoid purely definition-only questions.",
            "hard": "- Use scenario-based or calculation/logic-based stems.\n- Distractors should be very close and require precision.",
        },
        "tf": {
            "easy": "- Statements should be direct and explicit from the text.",
            "medium": "- Statements should involve interpretation, not exact sentence matching.",
            "hard": "- Statements should combine multiple facts; students must reason to judge truth.",
        },
        "fill": {
            "easy": "- Blank out high-salience keywords from direct definitions.",
            "medium": "- Blank out concept relationships or method names in applied contexts.",
            "hard": "- Blank out terms that require deep context understanding, not surface recall.",
        },
        "essay": {
            "easy": "- Ask for short explanations of core ideas.",
            "medium": "- Ask students to compare/contrast or apply a concept to a small scenario.",
            "hard": "- Ask for structured argument/analysis with justification and tradeoffs.",
        },
    }
    per_type = difficulty_by_type.get(question_type, {}).get(difficulty, "")
    return f"{base.get(question_type, 'Generate questions based on the content')}\n{per_type}".strip()


class QuizGenerateAPIView(APIView):
    permission_classes = [IsAuthenticated]
    
    def post(self, request):
        # Get quiz configuration from request
        question_type = str(request.data.get("questionType", "mcq")).lower()
        if question_type in {"true_false", "true/false"}:
            question_type = "tf"
        elif question_type in {"fill_in", "fill-in", "fillin"}:
            question_type = "fill"
        elif question_type in {"theory"}:
            question_type = "essay"
        try:
            question_count = int(request.data.get("questionCount", 10))
        except Exception:
            question_count = 10
        if question_count < 1:
            question_count = 1
        max_questions = max(5, int(getattr(settings, "QUIZ_MAX_QUESTIONS", 50)))
        if question_count > max_questions:
            question_count = max_questions
        difficulty = str(request.data.get("difficulty", "medium")).lower().strip()
        if difficulty not in {"easy", "medium", "hard"}:
            difficulty = "medium"
        chat_id = request.data.get("chat_id")

        profile = _get_payments_profile(request.user)
        premium_active = _is_premium(profile)
        if not premium_active and question_type in {"tf", "fill", "essay"}:
            return Response(
                {"error": "This question type requires a premium subscription."},
                status=403
            )
        if not premium_active and question_type != "mcq":
            remaining = 20 - profile.questions_generated
            if remaining <= 0 or question_count > remaining:
                return Response(
                    {"error": "You have reached your free question limit. Upgrade to premium to continue."},
                    status=429
                )
        
        # Get uploaded file
        if 'file' not in request.FILES:
            return Response({"error": "No file uploaded"}, status=400)
        
        uploaded_file = request.FILES['file']
        max_file_bytes = max(1024, int(getattr(settings, "QUIZ_SOURCE_MAX_FILE_BYTES", 50 * 1024 * 1024)))
        if uploaded_file.size > max_file_bytes:
            return Response(
                {"error": f"File too large. Maximum size is {max_file_bytes // (1024 * 1024)}MB."},
                status=400
            )
        
        # Validate file type
        allowed_types = ['pdf', 'txt', 'doc', 'docx', 'ppt', 'pptx', 'md']
        ext = uploaded_file.name.lower().split('.')[-1]
        if ext not in allowed_types:
            return Response(
                {"error": f"Invalid file type. Allowed: {', '.join(allowed_types)}"},
                status=400
            )

        source_sha = _sha256_uploaded_file(uploaded_file)
        max_chars = max(6000, int(getattr(settings, "QUIZ_SOURCE_MAX_CHARS", 20000)))
        max_pages = max(1, int(getattr(settings, "QUIZ_SOURCE_MAX_PAGES", 60)))
        quiz_cache_ttl = max(60, int(getattr(settings, "QUIZ_CACHE_TTL_SECONDS", 900)))
        quiz_hash_basis = json.dumps(
            {
                "user": request.user.id,
                "type": question_type,
                "difficulty": difficulty,
                "count": question_count,
                "source_sha": source_sha,
                "ext": ext,
                "max_chars": max_chars,
                "max_pages": max_pages,
            },
            sort_keys=True,
        )
        quiz_cache_key = f"quiz:v2:{hashlib.sha256(quiz_hash_basis.encode('utf-8')).hexdigest()}"
        normalized_questions = cache.get(quiz_cache_key)

        text_content = None
        if normalized_questions is None:
            context_cache_key = f"quizctx:v1:{hashlib.sha256(f'{source_sha}:{ext}:{max_chars}:{max_pages}'.encode('utf-8')).hexdigest()}"
            text_content = cache.get(context_cache_key)
            if text_content is None:
                text_content = _prepare_quiz_source_context(uploaded_file, max_chars, max_pages)
                if not text_content or len(text_content) < 100:
                    return Response(
                        {"error": "Could not extract enough content from the file. Please try a different file."},
                        status=400
                    )
                cache.set(context_cache_key, text_content, timeout=max(quiz_cache_ttl, 1800))
        
        # Generate quiz using OpenAI
        if normalized_questions is None:
            try:
                format_instructions = _quiz_format_instructions(question_type, difficulty)
                difficulty_instructions = _quiz_difficulty_instructions(difficulty)

                prompt = f"""
You are a study assistant. Based on the following content, generate {question_count} {difficulty} difficulty questions.

Content:
{text_content}

Question Type: {question_type}
Difficulty Rules:
{difficulty_instructions}

Question Format Rules:
{format_instructions}

Math Rendering Rules:
- If a question or explanation includes equations, format inline math as $...$.
- For multi-line derivations/equations, format block math as $$...$$.
- Do not emit raw LaTeX outside $...$ or $$...$$.
- Keep JSON escaping valid.

IMPORTANT: Return ONLY a valid JSON array with this EXACT structure:
[
    {{
        "id": "q1",
        "question": "Question text here?",
        "options": ["A", "B", "C", "D"],
        "correctAnswer": 0,
        "explanation": "Brief explanation",
        "concept": "Concept label for this question",
        "guidance": "Short study direction for this question",
        "correctText": "Expected answer text for fill/essay (optional for mcq/tf)"
    }}
]

For multiple choice: options must be an array of 4 strings, correctAnswer must be the index (0, 1, 2, or 3).
For true/false: options must be ["True", "False"], correctAnswer must be 0 or 1.
For fill/essay: options must be an empty array [], correctAnswer must be 0.

DO NOT include any markdown code blocks, DO NOT include any other text.
Only return valid JSON, nothing else.
"""
                client = get_openai_client()
                response = client.chat.completions.create(
                    model=getattr(settings, "QUIZ_OPENAI_MODEL", "gpt-4o-mini"),
                    messages=[
                        {"role": "system", "content": "You are a helpful study assistant that generates quizzes. Always return valid JSON only. Never use markdown code blocks."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.3,
                    max_tokens=2000
                )

                result = response.choices[0].message.content
                logger.info(f"[QuizGen] Raw OpenAI response (first 500 chars): {result[:500]}")

                questions = None
                parse_error = None
                try:
                    questions = json.loads(result)
                    logger.info(f"[QuizGen] Direct JSON parse succeeded, got {len(questions) if questions else 0} questions")
                except json.JSONDecodeError as e:
                    parse_error = f"Direct parse failed: {e}"
                    logger.warning(f"[QuizGen] {parse_error}")

                    json_match = re.search(r'```json?\n(.*?)\n```', result, re.DOTALL)
                    if json_match:
                        try:
                            questions = json.loads(json_match.group(1))
                            logger.info(f"[QuizGen] Markdown code block parse succeeded, got {len(questions)} questions")
                        except json.JSONDecodeError as e2:
                            parse_error = f"Markdown parse failed: {e2}"
                            logger.warning(f"[QuizGen] {parse_error}")
                    else:
                        logger.warning("[QuizGen] No markdown code block found")

                if questions is None:
                    start = result.find('[')
                    end = result.rfind(']') + 1
                    if start >= 0 and end > start:
                        try:
                            questions = json.loads(result[start:end])
                            logger.info(f"[QuizGen] Extracted JSON array from text, got {len(questions)} questions")
                        except json.JSONDecodeError as e3:
                            parse_error = f"Text extraction failed: {e3}"
                            logger.error(f"[QuizGen] {parse_error}")

                if questions is None:
                    logger.error(f"[QuizGen] All parsing methods failed. Last error: {parse_error}")
                    return Response(
                        {"error": "Failed to parse quiz questions. Please try again."},
                        status=500
                    )

                normalized_questions = []
                for i, q in enumerate(questions):
                    try:
                        if not isinstance(q, dict):
                            continue
                        normalized_q = {
                            "id": q.get("id", f"q{i+1}"),
                            "question": str(q.get("question", f"Question {i+1}")),
                            "options": [],
                            "correctAnswer": 0,
                            "explanation": str(q.get("explanation", "")),
                            "concept": str(q.get("concept", "")).strip(),
                            "guidance": str(q.get("guidance", "")).strip(),
                            "correctText": str(q.get("correctText", "")).strip(),
                        }
                        if not normalized_q["concept"]:
                            normalized_q["concept"] = _infer_question_concept(normalized_q["question"])
                        if not normalized_q["guidance"]:
                            normalized_q["guidance"] = "Focus on the concept behind this question and explain why the correct choice fits best."

                        if question_type == "mcq":
                            normalized_q["options"] = _normalize_quiz_options(q.get("options", []))
                            normalized_q["correctAnswer"] = _resolve_mcq_correct_index(
                                q.get("correctAnswer"),
                                normalized_q["options"],
                                normalized_q["correctText"],
                                normalized_q["explanation"],
                            )
                            if not normalized_q["correctText"] and normalized_q["options"]:
                                normalized_q["correctText"] = normalized_q["options"][normalized_q["correctAnswer"]]
                        elif question_type == "tf":
                            normalized_q["options"] = ["True", "False"]
                            ca = q.get("correctAnswer")
                            normalized_q["correctAnswer"] = 1 if ca in [1, "1", False, "False", "false"] else 0
                            if not normalized_q["correctText"]:
                                normalized_q["correctText"] = normalized_q["options"][normalized_q["correctAnswer"]]
                        else:
                            normalized_q["options"] = []
                            normalized_q["correctAnswer"] = 0
                            if not normalized_q["correctText"]:
                                normalized_q["correctText"] = normalized_q["explanation"]

                        normalized_questions.append(normalized_q)
                    except Exception as e:
                        logger.warning(f"[QuizGen] Failed to normalize question {i}: {e}")
                        continue

                logger.info(f"[QuizGen] Successfully normalized {len(normalized_questions)} questions")
                if not normalized_questions:
                    return Response(
                        {"error": "Failed to generate valid quiz questions. Please try again."},
                        status=500
                    )

                cache.set(quiz_cache_key, normalized_questions, timeout=quiz_cache_ttl)

            except Exception as e:
                logger.error(f"Quiz generation failed: {str(e)}", exc_info=True)
                return Response(
                    {"error": "Failed to generate quiz. Please try again."},
                    status=500
                )

        try:
            uploaded_file.seek(0)
        except Exception:
            pass
        persisted_source = _persist_source_upload(request.user, uploaded_file, purpose="qa")

        try:
            chat = _get_chat_for_write(request.user, chat_id)
            chat_title = None
            if chat:
                ChatHistory.objects.create(
                    chat=chat,
                    user=request.user,
                    prompt="",
                    response=f"I've generated {len(normalized_questions)} quiz questions based on your study material. Test your knowledge!",
                    prompt_type="text",
                    response_type="quiz",
                    response_metadata={
                        "title": "Generated Quiz",
                        "questions": normalized_questions,
                        "difficulty": difficulty,
                        "type": question_type,
                        "sourceFile": {
                            "id": persisted_source.id,
                            "name": persisted_source.file.name.rsplit("/", 1)[-1],
                            "storage_key": persisted_source.file.name,
                            "file_type": persisted_source.file_type,
                        },
                    },
                    context="quiz",
                )
                chat_title = _auto_title_feature_chat(
                    chat,
                    persisted_source.file.name.rsplit("/", 1)[-1],
                    "Quiz",
                )
            
            if not premium_active and question_type != "mcq":
                profile.questions_generated += len(normalized_questions)
                profile.save(update_fields=["questions_generated"])
            
            return Response({
                "questions": normalized_questions,
                "type": question_type,
                "difficulty": difficulty,
                "count": len(normalized_questions),
                "source_file": {
                    "id": persisted_source.id,
                    "name": persisted_source.file.name.rsplit("/", 1)[-1],
                    "storage_key": persisted_source.file.name,
                    "file_type": persisted_source.file_type,
                },
                "chat_title": chat_title,
            }, status=200)
            
        except Exception as e:
            logger.error(f"Quiz post-processing failed: {str(e)}", exc_info=True)
            return Response(
                {"error": "Failed to generate quiz. Please try again."},
                status=500
            )


class QuizGradeTextAPIView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        question = str(request.data.get("question", "")).strip()
        answer = str(request.data.get("answer", "")).strip()
        question_type = str(request.data.get("questionType", "essay")).lower().strip()
        difficulty = str(request.data.get("difficulty", "medium")).lower().strip()
        expected_answer = str(request.data.get("expectedAnswer", "")).strip()
        explanation = str(request.data.get("explanation", "")).strip()
        concept = str(request.data.get("concept", "")).strip()

        if question_type not in {"fill", "essay"}:
            return Response({"error": "Only fill/essay grading is supported."}, status=400)
        if not question or not answer:
            return Response({"error": "Question and answer are required."}, status=400)

        if difficulty not in {"easy", "medium", "hard"}:
            difficulty = "medium"

        prompt = f"""
You are grading a student answer for a study app.

Return ONLY valid JSON with this exact structure:
{{
  "is_correct": true,
  "score": 0.78,
  "concept": "short concept label",
  "guidance": "2-3 sentence actionable study direction",
  "feedback": "brief explanation of why this score was given"
}}

Rules:
- is_correct should be true only if the answer is reasonably correct for the question type/difficulty.
- score must be a float between 0 and 1.
- concept should be concise (2-5 words).
- guidance must help the student improve on this exact concept.
- feedback must be brief and specific.

Difficulty: {difficulty}
Question type: {question_type}
Question: {question}
Student answer: {answer}
Reference answer: {expected_answer}
Reference explanation: {explanation}
Known concept label: {concept}
"""
        try:
            client = get_openai_client()
            response = client.chat.completions.create(
                model=getattr(settings, "QUIZ_OPENAI_MODEL", "gpt-4o-mini"),
                messages=[
                    {"role": "system", "content": "You are a strict but supportive quiz grader. Return valid JSON only."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.1,
                max_tokens=350,
            )
            raw = (response.choices[0].message.content or "").strip()
            import json
            try:
                data = json.loads(raw)
            except json.JSONDecodeError:
                start = raw.find("{")
                end = raw.rfind("}") + 1
                data = json.loads(raw[start:end]) if start >= 0 and end > start else {}

            is_correct = bool(data.get("is_correct", False))
            try:
                score = float(data.get("score", 0.0))
            except Exception:
                score = 0.0
            score = max(0.0, min(1.0, score))
            concept_out = str(data.get("concept", concept or _infer_question_concept(question))).strip() or _infer_question_concept(question)
            guidance_out = str(data.get("guidance", "")).strip() or "Review the core concept, then rewrite your answer with one clear example."
            feedback_out = str(data.get("feedback", "")).strip() or "Answer graded."

            return Response({
                "is_correct": is_correct,
                "score": score,
                "concept": concept_out,
                "guidance": guidance_out,
                "feedback": feedback_out,
            }, status=200)
        except Exception as e:
            logger.error(f"Quiz text grading failed: {e}", exc_info=True)
            return Response({"error": "Failed to grade answer right now."}, status=500)

class SummarizeAPIView(APIView):
    permission_classes = [IsAuthenticated]
    
    def post(self, request):
        # Get summary options from request
        length = request.data.get("length", "detailed")  # brief, detailed, comprehensive
        format_type = request.data.get("format", "bullets")  # bullets, paragraphs
        include_key_terms = _to_bool(request.data.get("includeKeyTerms", True), default=True)
        chat_id = request.data.get("chat_id")
        
        # Get uploaded file
        if 'file' not in request.FILES:
            return Response({"error": "No file uploaded"}, status=400)
        
        uploaded_file = request.FILES['file']
        max_file_bytes = max(1024, int(getattr(settings, "SUMMARY_SOURCE_MAX_FILE_BYTES", 50 * 1024 * 1024)))
        if uploaded_file.size > max_file_bytes:
            return Response(
                {"error": f"File too large. Maximum size is {max_file_bytes // (1024 * 1024)}MB."},
                status=400
            )
        
        # Validate file type
        allowed_types = ['pdf', 'txt', 'doc', 'docx', 'ppt', 'pptx', 'md']
        ext = uploaded_file.name.lower().split('.')[-1]
        if ext not in allowed_types:
            return Response(
                {"error": f"Invalid file type. Allowed: {', '.join(allowed_types)}"},
                status=400
            )
        
        # Extract text from file
        text_content = extract_text_from_file(uploaded_file)
        
        min_chars = max(10, int(getattr(settings, "SUMMARY_SOURCE_MIN_CHARS", 20)))
        if not text_content or len(text_content.strip()) < min_chars:
            return Response(
                {"error": "Could not extract enough content from the file. Please try a different file or longer text."},
                status=400
            )
        
        # Determine word count based on length
        word_counts = {
            "brief": "~100 words",
            "detailed": "~300 words",
            "comprehensive": "~500 words"
        }
        
        max_chars = max(3000, int(getattr(settings, "SUMMARY_SOURCE_MAX_CHARS", 8500)))
        text_content = _prepare_llm_context(text_content, max_chars)

        summary_cache_ttl = max(60, int(getattr(settings, "SUMMARY_CACHE_TTL_SECONDS", 900)))
        summary_hash_basis = json.dumps(
            {
                "user": request.user.id,
                "length": str(length),
                "format": str(format_type),
                "include_key_terms": bool(include_key_terms),
                "content_sha": hashlib.sha256(text_content.encode("utf-8", errors="ignore")).hexdigest(),
            },
            sort_keys=True,
        )
        summary_cache_key = f"summary:v2:{hashlib.sha256(summary_hash_basis.encode('utf-8')).hexdigest()}"
        normalized_summary = cache.get(summary_cache_key)

        if normalized_summary is None:
            try:
                format_instruction = "Use bullet points" if format_type == "bullets" else "Use paragraphs"
                key_terms_instruction = "Include a section for key terms and definitions." if include_key_terms else ""

                prompt = f"""
You are a study assistant. Create a {length} summary of the following content.

Content:
{text_content}

Requirements:
- Length: {word_counts.get(length, '~300 words')}
- Format: {format_instruction}
- {key_terms_instruction if key_terms_instruction else "Do not include a key terms section."}
- Make it easy to study from
- Highlight important concepts
- Use clean Markdown for structure.
- If equations appear, use inline math as $...$ and block math as $$...$$.
- Do not output raw LaTeX fragments outside math delimiters.

Return the response as a JSON object with this structure:
{{
    "summary": "The summary content formatted as requested",
    "takeaways": ["Key point 1", "Key point 2", "Key point 3"],
    "keyTerms": [{{"term": "term name", "definition": "brief definition"}}]  // Only if includeKeyTerms is true
}}
Only return valid JSON, no additional text.
"""

                client = get_openai_client()
                response = client.chat.completions.create(
                    model=getattr(settings, "SUMMARY_OPENAI_MODEL", "gpt-4o-mini"),
                    messages=[
                        {"role": "system", "content": "You are a helpful study assistant that creates summaries. Always return valid JSON only."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.2,
                    max_tokens=1500
                )

                result = response.choices[0].message.content
                try:
                    summary_data = json.loads(result)
                except json.JSONDecodeError:
                    json_match = re.search(r'```json?\n(.*?)\n```', result, re.DOTALL)
                    if json_match:
                        summary_data = json.loads(json_match.group(1))
                    else:
                        start = result.find('{')
                        end = result.rfind('}') + 1
                        if start >= 0 and end > start:
                            summary_data = json.loads(result[start:end])
                        else:
                            raise ValueError("Could not parse summary from response")

                normalized_summary = _normalize_summary_payload(summary_data, include_key_terms)
                cache.set(summary_cache_key, normalized_summary, timeout=summary_cache_ttl)
            except Exception as e:
                logger.error(f"Summarization failed: {str(e)}", exc_info=True)
                return Response(
                    {"error": "Failed to generate summary. Please try again."},
                    status=500
                )

        try:
            uploaded_file.seek(0)
        except Exception:
            pass
        persisted_source = _persist_source_upload(request.user, uploaded_file, purpose="summary")

        try:
            summary_content = normalized_summary["summary"]
            chat = _get_chat_for_write(request.user, chat_id)
            chat_title = None
            if chat:
                ChatHistory.objects.create(
                    chat=chat,
                    user=request.user,
                    prompt="",
                    response=summary_content,
                    prompt_type="text",
                    response_type="summary",
                    response_metadata={
                        "title": "Document Summary",
                        "summary": summary_content,
                        "takeaways": normalized_summary["takeaways"],
                        "keyTerms": normalized_summary["keyTerms"],
                        "length": length,
                        "format": format_type,
                        "sourceFile": {
                            "id": persisted_source.id,
                            "name": persisted_source.file.name.rsplit("/", 1)[-1],
                            "storage_key": persisted_source.file.name,
                            "file_type": persisted_source.file_type,
                        },
                    },
                    context="summary",
                )
                chat_title = _auto_title_feature_chat(
                    chat,
                    persisted_source.file.name.rsplit("/", 1)[-1],
                    "Summary",
                )

            return Response({
                "summary": summary_content,
                "takeaways": normalized_summary["takeaways"],
                "keyTerms": normalized_summary["keyTerms"],
                "length": length,
                "format": format_type,
                "source_file": {
                    "id": persisted_source.id,
                    "name": persisted_source.file.name.rsplit("/", 1)[-1],
                    "storage_key": persisted_source.file.name,
                    "file_type": persisted_source.file_type,
                },
                "chat_title": chat_title,
            }, status=200)
        except Exception as e:
            logger.error(f"Summary post-processing failed: {str(e)}", exc_info=True)
            return Response(
                {"error": "Failed to generate summary. Please try again."},
                status=500
            )
